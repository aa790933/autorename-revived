from __future__ import annotations

import csv
import logging
import os
from abc import ABC, abstractmethod
from typing import Any, Dict, List

from autorename_revived._ai_processing import extract_metadata, extract_vision_metadata
from autorename_revived._pdf_utils import (
    ExtractionResult,
    assess_text_quality,
    extract_text_from_pdf,
    render_pages_to_images,
)

log = logging.getLogger(__name__)

_TEXT_QUALITY_THRESHOLD = 0.2

IMAGE_EXTENSIONS = frozenset({
    ".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp", ".gif", ".webp",
})


class BaseExtractor(ABC):
    def __init__(self, config: dict):
        self.config = config

    @abstractmethod
    def extract(self, filepath: str) -> ExtractionResult:
        ...

    @abstractmethod
    def extract_metadata(self, filepath: str) -> Dict[str, Any]:
        ...

    @staticmethod
    def _metadata_to_dict(md: Any) -> Dict[str, Any]:
        if hasattr(md, "model_dump"):
            return md.model_dump()
        if hasattr(md, "dict"):
            return md.dict()
        return dict(md)


class PdfExtractor(BaseExtractor):
    def extract(self, filepath: str) -> ExtractionResult:
        return extract_text_from_pdf(filepath)

    def extract_metadata(self, filepath: str) -> Dict[str, Any]:
        text_result = extract_text_from_pdf(filepath)
        quality = assess_text_quality(text_result.text)
        pdf_config = self.config.get("pdf", {})
        threshold = pdf_config.get("text_quality_threshold", _TEXT_QUALITY_THRESHOLD)

        if quality >= threshold:
            return self._metadata_to_dict(extract_metadata(text_result.text, self.config))

        vision = pdf_config.get("vision", "auto")
        if not vision:
            return self._metadata_to_dict(extract_metadata(text_result.text, self.config))

        if vision in ("auto", True):
            page_images = render_pages_to_images(filepath, max_pages=3)
            if page_images:
                import tempfile
                tmp_paths = []
                try:
                    for img_bytes in page_images:
                        tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
                        tmp.write(img_bytes)
                        tmp.close()
                        tmp_paths.append(tmp.name)
                    return self._metadata_to_dict(extract_vision_metadata(tmp_paths, self.config))
                except Exception:
                    if quality > 0:
                        return self._metadata_to_dict(extract_metadata(text_result.text, self.config))
                finally:
                    for tp in tmp_paths:
                        try:
                            os.unlink(tp)
                        except OSError:
                            pass

        return self._metadata_to_dict(extract_metadata(text_result.text, self.config))


class ImageExtractor(BaseExtractor):
    def extract(self, filepath: str) -> ExtractionResult:
        text_parts = []
        try:
            from PIL import Image
            from PIL.ExifTags import TAGS
            img = Image.open(filepath)
            exif_data = img.getexif()
            if exif_data:
                for tag_id, value in exif_data.items():
                    tag_name = TAGS.get(tag_id, str(tag_id))
                    if tag_name in ("ImageDescription", "XPTitle", "XPSubject", "XPComment"):
                        if value:
                            text_parts.append(str(value))
                    elif tag_name in ("DateTimeOriginal", "DateTime", "DateTimeDigitized"):
                        if value:
                            text_parts.append(str(value))
                    elif tag_name in ("Make", "Model"):
                        if value:
                            text_parts.append(str(value))
            img.close()
        except Exception as e:
            log.debug(f"EXIF extraction failed for {filepath}: {e}")
        return ExtractionResult(text="\n".join(text_parts), method="image-exif")

    def extract_metadata(self, filepath: str) -> Dict[str, Any]:
        try:
            return self._metadata_to_dict(extract_vision_metadata(filepath, self.config))
        except Exception as e:
            log.warning(f"Vision extraction failed for image {filepath}: {e}")
            result = self.extract(filepath)
            if result.text.strip():
                try:
                    return self._metadata_to_dict(extract_metadata(result.text, self.config))
                except Exception:
                    pass
            return {}


class DocxExtractor(BaseExtractor):
    def extract(self, filepath: str) -> ExtractionResult:
        try:
            from docx import Document
            doc = Document(filepath)
            text = "\n".join(p.text for p in doc.paragraphs)
            return ExtractionResult(text=text, method="docx")
        except Exception as e:
            return ExtractionResult(text="", method="docx", error=str(e))

    def extract_metadata(self, filepath: str) -> Dict[str, Any]:
        result = self.extract(filepath)
        if result.error:
            raise ValueError(f"DOCX extraction failed: {result.error}")
        return self._metadata_to_dict(extract_metadata(result.text, self.config))


class XlsxExtractor(BaseExtractor):
    def extract(self, filepath: str) -> ExtractionResult:
        try:
            from openpyxl import load_workbook
            wb = load_workbook(filepath, read_only=True, data_only=True)
            text_parts = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c is not None]
                    if cells:
                        text_parts.append(" ".join(cells))
            wb.close()
            return ExtractionResult(text="\n".join(text_parts), method="xlsx")
        except Exception as e:
            return ExtractionResult(text="", method="xlsx", error=str(e))

    def extract_metadata(self, filepath: str) -> Dict[str, Any]:
        result = self.extract(filepath)
        if result.error:
            raise ValueError(f"XLSX extraction failed: {result.error}")
        return self._metadata_to_dict(extract_metadata(result.text, self.config))


class PptxExtractor(BaseExtractor):
    def extract(self, filepath: str) -> ExtractionResult:
        try:
            from pptx import Presentation
            prs = Presentation(filepath)
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text.strip())
            return ExtractionResult(text="\n".join(text_parts), method="pptx")
        except Exception as e:
            return ExtractionResult(text="", method="pptx", error=str(e))

    def extract_metadata(self, filepath: str) -> Dict[str, Any]:
        result = self.extract(filepath)
        if result.error:
            raise ValueError(f"PPTX extraction failed: {result.error}")
        return self._metadata_to_dict(extract_metadata(result.text, self.config))


class CsvExtractor(BaseExtractor):
    def extract(self, filepath: str) -> ExtractionResult:
        try:
            text_parts = []
            for encoding in ("utf-8", "utf-8-sig", "latin-1", "cp1252"):
                try:
                    with open(filepath, "r", encoding=encoding, errors="strict") as f:
                        reader = csv.reader(f)
                        for i, row in enumerate(reader):
                            if i > 50:
                                break
                            cells = [c.strip() for c in row if c.strip()]
                            if cells:
                                text_parts.append(" ".join(cells))
                    break
                except (UnicodeDecodeError, UnicodeError):
                    text_parts.clear()
                    continue
            return ExtractionResult(text="\n".join(text_parts), method="csv")
        except Exception as e:
            return ExtractionResult(text="", method="csv", error=str(e))

    def extract_metadata(self, filepath: str) -> Dict[str, Any]:
        result = self.extract(filepath)
        if result.error:
            raise ValueError(f"CSV extraction failed: {result.error}")
        return self._metadata_to_dict(extract_metadata(result.text, self.config))


class TxtExtractor(BaseExtractor):
    def extract(self, filepath: str) -> ExtractionResult:
        try:
            for encoding in ("utf-8", "utf-8-sig", "latin-1", "cp1252"):
                try:
                    with open(filepath, "r", encoding=encoding, errors="strict") as f:
                        text = f.read(10000)
                    return ExtractionResult(text=text, method="text")
                except (UnicodeDecodeError, UnicodeError):
                    continue
            return ExtractionResult(text="", method="text", error="Could not decode file")
        except Exception as e:
            return ExtractionResult(text="", method="text", error=str(e))

    def extract_metadata(self, filepath: str) -> Dict[str, Any]:
        result = self.extract(filepath)
        if result.error:
            raise ValueError(f"Text extraction failed: {result.error}")
        return self._metadata_to_dict(extract_metadata(result.text, self.config))


def get_extractor(filepath: str, config: dict) -> BaseExtractor:
    ext = os.path.splitext(filepath)[1].lower()
    if ext == ".pdf":
        return PdfExtractor(config)
    if ext == ".docx":
        return DocxExtractor(config)
    if ext == ".doc":
        log.warning(f"Legacy .doc format: {filepath}. Convert to .docx for best results.")
        return DocxExtractor(config)
    if ext == ".xlsx":
        return XlsxExtractor(config)
    if ext == ".xls":
        log.warning(f"Legacy .xls format: {filepath}. Convert to .xlsx for best results.")
        return XlsxExtractor(config)
    if ext == ".pptx":
        return PptxExtractor(config)
    if ext == ".csv":
        return CsvExtractor(config)
    if ext in (".txt", ".md", ".rtf"):
        return TxtExtractor(config)
    if ext in IMAGE_EXTENSIONS:
        return ImageExtractor(config)
    raise ValueError(f"Unsupported file type: {ext}")
